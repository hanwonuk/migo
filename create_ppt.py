#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
음원 마케팅 체험단 PPT 생성 스크립트
HTML 웹페이지 내용을 바탕으로 PowerPoint 프레젠테이션을 생성합니다.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
import os

def create_music_marketing_ppt():
    """음원 마케팅 체험단 PPT 생성"""
    
    # 새 프레젠테이션 생성
    prs = Presentation()
    
    # 색상 정의 (CSS 변수에서 가져온 색상)
    primary_color = RGBColor(99, 102, 241)    # #6366f1
    secondary_color = RGBColor(16, 185, 129)  # #10b981
    accent_color = RGBColor(245, 158, 11)     # #f59e0b
    dark_color = RGBColor(31, 41, 55)         # #1f2937
    light_color = RGBColor(107, 114, 128)     # #6b7280
    
    # 슬라이드 1: 타이틀 페이지
    create_title_slide(prs, primary_color, accent_color)
    
    # 슬라이드 2: 체험단 진행과정
    create_process_slide(prs, primary_color, secondary_color)
    
    # 슬라이드 3: 인스타그램/틱톡 프로그램
    create_instagram_tiktok_slide(prs, primary_color, accent_color)
    
    # 슬라이드 4: 음원 다운로드 프로그램
    create_music_download_slide(prs, primary_color, secondary_color)
    
    # 슬라이드 5: 블로그/카페/커뮤니티 프로그램
    create_blog_community_slide(prs, primary_color, accent_color)
    
    # 슬라이드 6: 체험단 특별 혜택
    create_benefits_slide(prs, primary_color, secondary_color)
    
    # 슬라이드 7: 신청 방법
    create_application_slide(prs, primary_color, accent_color)
    
    # PPT 파일 저장
    filename = '음원_마케팅_체험단.pptx'
    prs.save(filename)
    print(f"✅ PPT 파일이 생성되었습니다: {filename}")
    
    return filename

def create_title_slide(prs, primary_color, accent_color):
    """타이틀 슬라이드 생성"""
    slide_layout = prs.slide_layouts[6]  # 빈 슬라이드
    slide = prs.slides.add_slide(slide_layout)
    
    # 배경색 설정
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(103, 126, 234)  # 그라데이션 느낌의 보라색
    
    # 메인 제목
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "음원 마케팅 체험단"
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(48)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    
    # 부제목
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "새로운 음원의 시작, 미교 체험단"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.alignment = PP_ALIGN.CENTER
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.color.rgb = accent_color
    
    # 플로팅 요소들 (텍스트로 표현)
    elements = [
        ("📱 인스타그램", Inches(1.5), Inches(6)),
        ("🎬 틱톡", Inches(4.25), Inches(6.5)),
        ("🔍 네이버", Inches(7), Inches(6))
    ]
    
    for text, left, top in elements:
        element_box = slide.shapes.add_textbox(left, top, Inches(1.5), Inches(0.8))
        element_frame = element_box.text_frame
        element_frame.text = text
        element_para = element_frame.paragraphs[0]
        element_para.alignment = PP_ALIGN.CENTER
        element_para.font.size = Pt(14)
        element_para.font.color.rgb = RGBColor(255, 255, 255)

def create_process_slide(prs, primary_color, secondary_color):
    """체험단 진행과정 슬라이드"""
    slide_layout = prs.slide_layouts[5]  # 제목과 내용
    slide = prs.slides.add_slide(slide_layout)
    
    # 제목
    title = slide.shapes.title
    title.text = "체험단 진행과정"
    title_para = title.text_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.color.rgb = primary_color
    
    # 부제목
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(0.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "체험단 신청부터 완료까지의 단계별 진행과정을 안내해드립니다."
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.alignment = PP_ALIGN.CENTER
    subtitle_para.font.size = Pt(16)
    
    # 진행과정 단계들
    steps = [
        "1. 체험단 신청서 작성\n(신청자)",
        "2. 접수\n(대행사)",
        "3. 체험단 승인\n(대행사)",
        "4. 게시물 작성\n(신청자)",
        "5. 게시물 확인\n(대행사)",
        "6. 이체(입금)\n(대행사)"
    ]
    
    # 3x2 그리드로 배치
    for i, step in enumerate(steps):
        row = i // 3
        col = i % 3
        left = Inches(1 + col * 2.7)
        top = Inches(3 + row * 2)
        
        # 단계 박스 생성
        step_box = slide.shapes.add_textbox(left, top, Inches(2.5), Inches(1.5))
        step_frame = step_box.text_frame
        step_frame.text = step
        step_para = step_frame.paragraphs[0]
        step_para.alignment = PP_ALIGN.CENTER
        step_para.font.size = Pt(14)
        step_para.font.bold = True
        
        # 화살표 추가 (마지막 단계 제외)
        if i < len(steps) - 1 and col < 2:
            arrow_box = slide.shapes.add_textbox(left + Inches(2.5), top + Inches(0.5), Inches(0.2), Inches(0.5))
            arrow_frame = arrow_box.text_frame
            arrow_frame.text = "→"
            arrow_para = arrow_frame.paragraphs[0]
            arrow_para.font.size = Pt(20)
            arrow_para.font.color.rgb = primary_color

def create_instagram_tiktok_slide(prs, primary_color, accent_color):
    """인스타그램/틱톡 프로그램 슬라이드"""
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    
    # 제목
    title = slide.shapes.title
    title.text = "📱 인스타그램 / 틱톡"
    title_para = title.text_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.color.rgb = primary_color
    
    # 테이블 생성
    table_data = [
        ["방식", "구분", "과정", "추가 내용", "금액"],
        ["1", "인스타그램 업로드(릴스)", "본인 사진 및 영상 업로드", "게시물 업로드 시 지정된 노래 삽입", "3천원"],
        ["2", "인스타그램 업로드(릴스)", "대행사에서 사진 및 영상 전달", "게시물 업로드 시 지정된 노래 삽입", "3천원"],
        ["3", "틱톡 업로드", "대행사에서 사진 및 영상 전달", "게시물 업로드 시 지정된 노래 삽입", "4천원"]
    ]
    
    create_table_on_slide(slide, table_data, Inches(0.5), Inches(2.5), primary_color, accent_color)
    
    # 참고사항
    note_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.5))
    note_frame = note_box.text_frame
    note_frame.text = "※ 참고사항: 개인에 따라 약간의 편집이 필요함"
    note_para = note_frame.paragraphs[0]
    note_para.font.size = Pt(12)
    note_para.font.italic = True

def create_music_download_slide(prs, primary_color, secondary_color):
    """음원 다운로드 프로그램 슬라이드"""
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    
    # 제목
    title = slide.shapes.title
    title.text = "🎵 음원 다운로드"
    title_para = title.text_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.color.rgb = primary_color
    
    # 테이블 생성
    table_data = [
        ["방식", "구분", "과정", "추가 내용", "비고"],
        ["1", "멜론 음원 다운로드", "가수검색, 음악검색 → 음원 다운로드", "음원 듣기 및 공감", "3천원"],
        ["2", "네이버 검색", "해당 가수 또는 음원 검색 → 음원 플레이", "해당 게시물 클릭 및 리딩", "3천원"]
    ]
    
    create_table_on_slide(slide, table_data, Inches(0.5), Inches(2.5), primary_color, secondary_color)

def create_blog_community_slide(prs, primary_color, accent_color):
    """블로그/카페/커뮤니티 프로그램 슬라이드"""
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    
    # 제목
    title = slide.shapes.title
    title.text = "✍️ 블로그 / 카페 / 커뮤니티 후기 작성"
    title_para = title.text_frame.paragraphs[0]
    title_para.font.size = Pt(28)
    title_para.font.color.rgb = primary_color
    
    # 테이블 생성
    table_data = [
        ["방식", "구분", "과정", "추가 내용", "비고"],
        ["1", "블로그 후기작성", "음원다운로드 후 블로그 후기 작성", "감상평/후기글 1500자 이상 작성", "1.5만"],
        ["2", "카페글 작성", "음원다운로드 후 네이버 카페글 작성", "감상평/후기글 500자 이상 작성", "1만"],
        ["3", "커뮤니티글 작성", "음원다운로드 후 커뮤니티 글 게시", "감상평/후기글 500자 이상 작성", "1만"]
    ]
    
    create_table_on_slide(slide, table_data, Inches(0.5), Inches(2.5), primary_color, accent_color)

def create_benefits_slide(prs, primary_color, secondary_color):
    """체험단 특별 혜택 슬라이드"""
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    
    # 제목
    title = slide.shapes.title
    title.text = "🎁 체험단 특별 혜택"
    title_para = title.text_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.color.rgb = primary_color
    
    # 부제목
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(0.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "미교 체험단 참여자만의 특별한 혜택을 만나보세요."
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.alignment = PP_ALIGN.CENTER
    subtitle_para.font.size = Pt(16)
    
    # 혜택 박스들 (2x2 그리드)
    benefits = [
        ("🎁 무료 체험", "모든 프로그램을\n무료로 체험할 수 있습니다."),
        ("📜 수료증 발급", "프로그램 완료 시\n공식 수료증을 발급합니다."),
        ("💰 할인 혜택", "정규 프로그램 등록 시\n최대 50% 할인 혜택"),
        ("👥 커뮤니티", "체험단 전용 커뮤니티에서\n지속적인 교류")
    ]
    
    for i, (title_text, desc_text) in enumerate(benefits):
        row = i // 2
        col = i % 2
        left = Inches(1.5 + col * 4)
        top = Inches(3 + row * 2.5)
        
        # 혜택 박스
        benefit_box = slide.shapes.add_textbox(left, top, Inches(3.5), Inches(2))
        benefit_frame = benefit_box.text_frame
        
        # 제목 추가
        title_p = benefit_frame.paragraphs[0]
        title_p.text = title_text
        title_p.font.size = Pt(18)
        title_p.font.bold = True
        title_p.font.color.rgb = primary_color
        title_p.alignment = PP_ALIGN.CENTER
        
        # 설명 추가
        desc_p = benefit_frame.add_paragraph()
        desc_p.text = desc_text
        desc_p.font.size = Pt(14)
        desc_p.alignment = PP_ALIGN.CENTER

def create_application_slide(prs, primary_color, accent_color):
    """신청 방법 슬라이드"""
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    
    # 제목
    title = slide.shapes.title
    title.text = "📝 지금 바로 신청하세요!"
    title_para = title.text_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.color.rgb = primary_color
    
    # 부제목
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(0.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "미교 체험단과 함께 새로운 음원 마케팅의 미래를 경험해보세요."
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.alignment = PP_ALIGN.CENTER
    subtitle_para.font.size = Pt(16)
    
    # 신청 폼 필드들
    form_fields = [
        "✓ 이름",
        "✓ 이메일 주소",
        "✓ 연락처",
        "✓ 관심 프로그램 선택",
        "✓ 궁금한 점이나 요청사항"
    ]
    
    form_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(4), Inches(3))
    form_frame = form_box.text_frame
    form_frame.text = "신청 양식:\n\n" + "\n".join(form_fields)
    
    # 혜택 포인트
    benefits_text = """신청 혜택:

• 전문 강사진과의 1:1 상담
• 개인 맞춤형 학습 계획
• 24시간 학습 지원 서비스"""
    
    benefits_box = slide.shapes.add_textbox(Inches(5.5), Inches(2.5), Inches(4), Inches(3))
    benefits_frame = benefits_box.text_frame
    benefits_frame.text = benefits_text

def create_table_on_slide(slide, table_data, left, top, primary_color, accent_color):
    """슬라이드에 테이블 생성"""
    rows = len(table_data)
    cols = len(table_data[0])
    
    # 테이블 크기 계산
    table_width = Inches(9)
    table_height = Inches(3)
    
    table = slide.shapes.add_table(rows, cols, left, top, table_width, table_height).table
    
    # 테이블 데이터 입력
    for row_idx, row_data in enumerate(table_data):
        for col_idx, cell_data in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(cell_data)
            
            # 헤더 행 스타일링
            if row_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = primary_color
                paragraph = cell.text_frame.paragraphs[0]
                paragraph.font.color.rgb = RGBColor(255, 255, 255)
                paragraph.font.bold = True
                paragraph.font.size = Pt(12)
            else:
                # 데이터 행 스타일링
                paragraph = cell.text_frame.paragraphs[0]
                paragraph.font.size = Pt(10)
                
                # 금액 열 강조
                if col_idx == cols - 1 and row_idx > 0:
                    paragraph.font.color.rgb = accent_color
                    paragraph.font.bold = True

if __name__ == "__main__":
    print("🎬 음원 마케팅 체험단 PPT 생성을 시작합니다...")
    
    try:
        filename = create_music_marketing_ppt()
        print(f"🎉 성공적으로 PPT 파일이 생성되었습니다!")
        print(f"📁 파일 위치: {os.path.abspath(filename)}")
        
    except ImportError as e:
        print("❌ 오류: python-pptx 라이브러리가 설치되지 않았습니다.")
        print("💡 다음 명령어를 실행해주세요: pip install python-pptx")
        
    except Exception as e:
        print(f"❌ PPT 생성 중 오류가 발생했습니다: {e}") 